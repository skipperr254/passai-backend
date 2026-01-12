from pptx import Presentation
from io import BytesIO

def extract_pptx_text(bytes_data: bytes) -> str:
    presentation = Presentation(BytesIO(bytes_data))
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    # Remove null bytes that PostgreSQL cannot store
    text = text.replace('\x00', '')
    return text.strip()